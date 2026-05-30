"""
pptx_gyarto.py -- Marp Markdown → PPTX converter
Converts a Marp-formatted Markdown file to a PowerPoint presentation.

Usage:
    python scripts/pptx_gyarto.py <marp_md_file> [--template <template.pptx>] [--output <out.pptx>]

If --template is not provided (or file missing), uses a built-in default style.

NOTE (pipeline): du_template.pptx is the intended template. If missing,
this script generates a functionally equivalent default-styled deck.
"""

import argparse
import re
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    sys.exit("python-pptx not installed. Run: pip install python-pptx --break-system-packages")


# ---------------------------------------------------------------------------
# Slide dimensions: 16:9 widescreen
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Colours (neutral academic theme)
C_TITLE_BG  = RGBColor(0x1F, 0x49, 0x7D)   # dark blue
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK     = RGBColor(0x1A, 0x1A, 0x1A)
C_ACCENT    = RGBColor(0x2E, 0x75, 0xB6)   # medium blue
C_LIGHT_BG  = RGBColor(0xF2, 0xF2, 0xF2)   # light grey
C_MSC_BG    = RGBColor(0xFF, 0xF0, 0xD0)   # warm amber for MSc slides


def parse_marp(md_text: str) -> list[dict]:
    """
    Parse Marp Markdown into a list of slide dicts.
    Each dict: {title, subtitle, body, is_msc, is_cover}
    """
    # Strip YAML frontmatter
    md_text = re.sub(r'^---\n.*?---\n', '', md_text, count=1, flags=re.DOTALL)

    # Split on --- slide separators
    raw_slides = re.split(r'\n---\n', md_text)

    slides = []
    for raw in raw_slides:
        raw = raw.strip()
        if not raw:
            continue

        is_msc = bool(re.search(r'<!-- MSc -->', raw))
        # Strip MSc comments
        raw = re.sub(r'<!-- /?MSc -->', '', raw).strip()

        lines = raw.split('\n')
        title = ''
        subtitle = ''
        body_lines = []
        skip_next = False

        for i, line in enumerate(lines):
            if skip_next:
                skip_next = False
                continue
            if re.match(r'^# ', line):
                title = line[2:].strip()
                # Next line may be subtitle (###)
                if i + 1 < len(lines) and re.match(r'^### ', lines[i+1]):
                    subtitle = lines[i+1][4:].strip()
                    skip_next = True
            elif re.match(r'^## ', line):
                title = line[3:].strip()
            elif re.match(r'^### ', line):
                subtitle = line[4:].strip()
            else:
                body_lines.append(line)

        body = '\n'.join(body_lines).strip()
        slides.append({
            'title': title,
            'subtitle': subtitle,
            'body': body,
            'is_msc': is_msc,
            'is_cover': bool(not subtitle and re.match(r'^# ', raw.split('\n')[0])) if lines else False,
        })
    return slides


def add_textbox(slide, left, top, width, height, text, font_size=Pt(18),
                bold=False, color=None, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def fill_shape(shape, color: RGBColor):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def parse_body_segments(body_text: str) -> list[dict]:
    """
    Split body text into segments: text, image, or table.
    Returns list of {'type': str, 'content': str | list}.
    """
    segments = []
    current_text = []

    def flush_text():
        t = '\n'.join(current_text).strip()
        if t:
            segments.append({'type': 'text', 'content': t})
        current_text.clear()

    lines = body_text.split('\n')
    i = 0
    while i < len(lines):
        line = lines[i]

        # Image: ![alt](path)
        img_m = re.match(r'^\s*!\[([^\]]*)\]\(([^)]+)\)\s*$', line)
        if img_m:
            flush_text()
            segments.append({'type': 'image',
                             'alt': img_m.group(1),
                             'path': img_m.group(2)})
            i += 1
            continue

        # GFM table: line starts with |
        if re.match(r'^\s*\|', line):
            flush_text()
            table_lines = []
            while i < len(lines) and re.match(r'^\s*\|', lines[i]):
                table_lines.append(lines[i])
                i += 1
            segments.append({'type': 'table', 'content': table_lines})
            continue

        current_text.append(line)
        i += 1

    flush_text()
    return segments


def parse_gfm_table(table_lines: list[str]) -> tuple[list[str], list[list[str]]]:
    """Parse GFM table lines into (headers, rows). Skips separator row."""
    headers = []
    rows = []
    for line in table_lines:
        if re.match(r'^\s*\|[-: |]+\|\s*$', line):
            continue  # separator
        cells = [c.strip() for c in re.split(r'(?<!\\)\|', line) if c.strip() != '']
        if not headers:
            headers = cells
        else:
            rows.append(cells)
    return headers, rows


def add_pptx_table(slide, left, top, width, height,
                   headers: list[str], rows: list[list[str]],
                   font_size=Pt(13), header_color=None, text_color=None):
    """Add a python-pptx table to a slide."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    if header_color is None:
        header_color = RGBColor(0x1F, 0x49, 0x7D)
    if text_color is None:
        text_color = RGBColor(0x1A, 0x1A, 0x1A)

    n_rows = len(rows) + 1  # +1 for header
    n_cols = max(len(headers), max((len(r) for r in rows), default=0))
    if n_rows < 1 or n_cols < 1:
        return

    # Clamp height so table fits
    row_h = min(height // n_rows, Inches(0.45))

    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_h * n_rows).table

    def set_cell(tbl, row_idx, col_idx, text, bold=False, fg=None, bg=None):
        cell = tbl.cell(row_idx, col_idx)
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
            run.text = text
        run.font.size = font_size
        run.font.bold = bold
        if fg:
            run.font.color.rgb = fg
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    # Header row
    for c, hdr in enumerate(headers[:n_cols]):
        set_cell(tbl, 0, c, hdr, bold=True,
                 fg=RGBColor(0xFF, 0xFF, 0xFF), bg=header_color)

    # Data rows
    for r, row in enumerate(rows):
        for c in range(n_cols):
            val = row[c] if c < len(row) else ''
            set_cell(tbl, r + 1, c, val, fg=text_color)

    return tbl


def add_slide_image(slide, left, top, width, height, img_path: str, alt: str,
                    md_dir: str | None = None):
    """Add a picture to a slide. img_path may be relative to md_dir."""
    from pptx.util import Pt
    path = Path(img_path)
    if not path.is_absolute() and md_dir:
        path = Path(md_dir) / path
    if not path.exists():
        print(f"  ⚠️  Kép nem található, kihagyva: {path}")
        return None
    try:
        pic = slide.shapes.add_picture(str(path), left, top, width, height)
        return pic
    except Exception as e:
        print(f"  ⚠️  Kép betöltési hiba ({path}): {e}")
        return None


def build_presentation(slides_data: list[dict], template_path: str | None,
                       md_dir: str | None = None) -> Presentation:
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
        print(f"Template betöltve: {template_path}")
    else:
        prs = Presentation()
        prs.slide_width  = SLIDE_W
        prs.slide_height = SLIDE_H
        if template_path:
            print(f"⚠️  Helyettesítés: {template_path!r} nem található -- "
                  "du_template.pptx hiányában default stílust használok. "
                  "Dokumentáció: kepek_workflow.md §7 Teszt helyettesítés.")

    blank_layout = prs.slide_layouts[6]  # blank

    for idx, sd in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        w, h = prs.slide_width, prs.slide_height
        m = Inches(0.4)   # margin

        # Background colour
        bg_color = C_MSC_BG if sd['is_msc'] else (C_TITLE_BG if idx == 0 else C_WHITE)
        bg = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            0, 0, w, h
        )
        fill_shape(bg, bg_color)
        bg.line.fill.background()

        text_color = C_WHITE if (idx == 0 and not sd['is_msc']) else C_BLACK

        # Title bar (coloured strip at top for non-cover slides)
        if idx > 0:
            bar = slide.shapes.add_shape(1, 0, 0, w, Inches(1.1))
            fill_shape(bar, C_TITLE_BG)
            bar.line.fill.background()

        # Title
        if sd['title']:
            title_top = Inches(0.15) if idx == 0 else Inches(0.1)
            title_h   = Inches(1.0)
            add_textbox(
                slide, m, title_top, w - 2*m, title_h,
                sd['title'],
                font_size=Pt(32) if idx == 0 else Pt(26),
                bold=True,
                color=C_WHITE,
                align=PP_ALIGN.CENTER if idx == 0 else PP_ALIGN.LEFT,
            )

        # Subtitle (cover slide)
        if sd['subtitle'] and idx == 0:
            add_textbox(
                slide, m, Inches(1.4), w - 2*m, Inches(0.6),
                sd['subtitle'],
                font_size=Pt(20), bold=False,
                color=C_WHITE, align=PP_ALIGN.CENTER,
            )

        # MSc badge
        if sd['is_msc']:
            badge = slide.shapes.add_shape(1, w - Inches(1.8), Inches(0.05), Inches(1.7), Inches(0.45))
            fill_shape(badge, RGBColor(0xD0, 0x70, 0x00))
            badge.line.fill.background()
            add_textbox(slide, w - Inches(1.8), Inches(0.06), Inches(1.7), Inches(0.4),
                        "MSc", font_size=Pt(14), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Body
        if sd['body']:
            body_top = Inches(1.25) if idx > 0 else Inches(2.1)
            body_h   = h - body_top - Inches(0.3)

            # Pre-clean: code fences + HTML comments
            body_text = sd['body']
            body_text = re.sub(r'```[^\n]*\n', '', body_text)
            body_text = re.sub(r'```', '', body_text)
            body_text = re.sub(r'<!--.*?-->', '', body_text, flags=re.DOTALL)
            body_text = body_text.strip()

            # Parse into segments (text / image / table)
            segments = parse_body_segments(body_text)

            # Simple layout: divide body_h equally among segments
            n_seg = len(segments)
            seg_h = body_h // n_seg if n_seg else body_h
            cur_top = body_top

            for seg in segments:
                if seg['type'] == 'image':
                    # Leave room for alt-text caption below image
                    img_h = int(seg_h * 0.85)
                    add_slide_image(
                        slide, m, cur_top, w - 2*m, img_h,
                        seg['path'], seg.get('alt', ''),
                        md_dir=md_dir,
                    )
                    # Alt text below image (small)
                    if seg.get('alt'):
                        cap_top = cur_top + img_h
                        add_textbox(
                            slide, m, cap_top, w - 2*m, seg_h - img_h,
                            seg['alt'], font_size=Pt(11),
                            color=C_ACCENT, align=PP_ALIGN.CENTER,
                        )

                elif seg['type'] == 'table':
                    headers, rows = parse_gfm_table(seg['content'])
                    if headers:
                        add_pptx_table(
                            slide, m, cur_top, w - 2*m, seg_h,
                            headers, rows,
                            font_size=Pt(12),
                            text_color=text_color,
                        )

                else:  # text
                    t = seg['content'].strip()
                    if t:
                        n_chars = len(t)
                        fsize = Pt(13) if n_chars > 600 else Pt(15) if n_chars > 300 else Pt(17)
                        add_textbox(
                            slide, m, cur_top, w - 2*m, seg_h,
                            t, font_size=fsize, color=text_color,
                        )

                cur_top += seg_h

        # Page number (not on cover)
        if idx > 0:
            add_textbox(
                slide, w - Inches(1.2), h - Inches(0.35), Inches(1.0), Inches(0.3),
                str(idx),
                font_size=Pt(11), color=C_ACCENT, align=PP_ALIGN.RIGHT,
            )

    return prs


def _detect_week_number(week_dir: Path) -> int:
    """'1_het' → 1, '12_het' → 12"""
    import re
    m = re.match(r'^(\d+)_het$', week_dir.name)
    if m:
        return int(m.group(1))
    raise ValueError(f"Nem tudja meghatározni a hét számát: {week_dir.name!r}. "
                     "Elvárt formátum: N_het (pl. 1_het, 12_het).")


def main():
    parser = argparse.ArgumentParser(
        description='Marp MD → DUE PPTX (pipeline 12. lépés)',
        formatter_class=argparse.RawTextHelpFormatter,
        epilog=(
            "Pipeline-módban:\n"
            "  python 12_pptx_gyarto.py --week-dir test_outputs/meta_file_updates_test/1_het\n"
            "  Input:  4_wip_outputs/N_Prezentacio.md\n"
            "  Output: 5_clean_outputs/N_Prezentacio.pptx\n\n"
            "Közvetlen módban:\n"
            "  python 12_pptx_gyarto.py 1_Prezentacio.md --output out.pptx"
        )
    )
    parser.add_argument('input', nargs='?', default=None,
                        help='Marp Markdown fájl (.md) — elhagyható ha --week-dir megadva')
    parser.add_argument('--week-dir', default=None,
                        help='Pipeline hét-mappa (pl. test_outputs/<Tantargy>/N_het)')
    parser.add_argument('--template', default=None,
                        help='PPTX template elérési útja (default: templates/due_refactored.pptx)')
    parser.add_argument('--output', default=None,
                        help='Kimeneti .pptx fájl (elhagyható ha --week-dir megadva)')
    parser.add_argument('--pdf', action='store_true',
                        help='PDF másolat is generálódjon (PowerPoint COM, Windows only)')
    args = parser.parse_args()

    # -- Útvonalak meghatározása --
    default_template = Path('templates/due_refactored.pptx')

    if args.week_dir:
        week_dir = Path(args.week_dir)
        if not week_dir.is_dir():
            sys.exit(f"Nem található hét-mappa: {week_dir}")
        n = _detect_week_number(week_dir)
        md_path  = week_dir / '4_wip_outputs' / f'{n}_Prezentacio.md'
        out_path = week_dir / '5_clean_outputs' / f'{n}_Prezentacio.pptx'
        if args.input:
            md_path = Path(args.input)
        if args.output:
            out_path = Path(args.output)
    elif args.input:
        md_path  = Path(args.input)
        out_path = Path(args.output) if args.output else md_path.with_suffix('.pptx')
    else:
        parser.error("Kötelező: --week-dir VAGY positional input (md fájl)")

    if not md_path.exists():
        sys.exit(f"Nem található: {md_path}")

    template_path = args.template or str(default_template)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    with open(md_path, encoding='utf-8') as f:
        md_text = f.read()

    slides_data = parse_marp(md_text)
    print(f"Feldolgozva: {len(slides_data)} dia <- {md_path}")

    prs = build_presentation(slides_data, template_path, md_dir=str(md_path.parent))
    prs.save(out_path)
    print(f"Mentve: {out_path}")

    if args.pdf:
        _pptx_to_pdf(out_path)


def _pptx_to_pdf(pptx_path: Path):
    """PPTX → PDF via PowerPoint COM (Windows only)."""
    pdf_path = pptx_path.with_suffix('.pdf')
    try:
        import win32com.client
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        prs = ppt.Presentations.Open(str(pptx_path.resolve()), True, False, False)
        prs.SaveAs(str(pdf_path.resolve()), 32)  # 32 = ppSaveAsPDF
        prs.Close()
        ppt.Quit()
        print(f"PDF: {pdf_path}")
    except ImportError:
        print("  WARN  pywin32 nem elérhető -- PDF kihagyva (pip install pywin32)")
    except Exception as e:
        print(f"  WARN  PDF generálás sikertelen: {e}")


if __name__ == '__main__':
    main()
