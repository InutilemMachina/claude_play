"""
due_mindmap_fill.py  --  DUE mindmap breadcrumb automatizmus
============================================================

A mindmap-variáns lényege: minden content dián egy jobb oldali "mindmap"
sidebar mutatja az adott dia helyét a fejezet-hierarchiában (breadcrumb).

Ez a modul AUTOMATIKUSAN generálja ezt a breadcrumb-ot egy fejezet-fából,
és feltölti a mindmap template (.pptx) `mindmap_body` shape-jeit a megfelelő
útvonallal — diánként a dia logikai pozíciója alapján.

Fő API:
    tree = ChapterTree([
        Chapter("1", "Első fejezet", [
            Section("1.1", "Első szakasz", ["1.1.1 Első alszakasz"]),
        ]),
        Chapter("2", "Második fejezet"),
    ])

    fill_mindmap(
        "templates/due_prenetation_template_mindmap.pptx",
        "output/eloadas.pptx",
        tree,
        # diánkénti pozíció: slide-index -> a fejezet-fa csomópont azonosítója
        positions={
            2: "1",      # 3. dia (0-indexelt 2) az "1." fejezetnél tart
            3: "1.1.1",  # ...
        },
    )

A breadcrumb a gyökértől az aktuális csomópontig terjedő utat mutatja,
az aktuális csomópontot kiemelve (bold + narancs háttér-jelölés).
"""

from __future__ import annotations
import copy
from dataclasses import dataclass, field
from pathlib import Path

try:
    from pptx import Presentation
    from lxml import etree
except ImportError:
    raise ImportError("pip install python-pptx lxml")

# ---------------------------------------------------------------------------
# Színek / formázás (a mindmap template méréseiből)
# ---------------------------------------------------------------------------
C_ORANGE = "D4622A"      # sorszám prefix
C_DARK   = "1A1A2E"      # szöveg
C_CURRENT_BG = "FCE9DF"  # aktuális csomópont halvány narancs háttér
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ---------------------------------------------------------------------------
# Fejezet-fa adatszerkezet
# ---------------------------------------------------------------------------

@dataclass
class Node:
    """Egy csomópont a fejezet-fában (fejezet / szakasz / alszakasz)."""
    number: str                       # pl. "1", "1.1", "1.1.1"
    title: str                        # pl. "Első fejezet"
    children: list["Node"] = field(default_factory=list)

    @property
    def level(self) -> int:
        """0 = fejezet, 1 = szakasz, 2 = alszakasz (a number pontjaiból)."""
        return self.number.count(".")


def Chapter(number, title, sections=None):
    """Kényelmi konstruktor: fejezet, opcionális szakaszokkal."""
    return Node(number, title, sections or [])


def Section(number, title, subs=None):
    """Kényelmi konstruktor: szakasz, opcionális alszakaszokkal (str vagy Node)."""
    children = []
    for s in (subs or []):
        if isinstance(s, Node):
            children.append(s)
        else:
            # "1.1.1 Cím" formátumú string
            num, _, ttl = s.partition(" ")
            children.append(Node(num, ttl))
    return Node(number, title, children)


class ChapterTree:
    """A teljes fejezet-hierarchia + útvonal-kereső."""

    def __init__(self, chapters: list[Node]):
        self.chapters = chapters
        self._index: dict[str, list[Node]] = {}
        self._build_index(chapters, [])

    def _build_index(self, nodes, ancestors):
        for n in nodes:
            path = ancestors + [n]
            self._index[n.number] = path
            self._build_index(n.children, path)

    def path_to(self, number: str) -> list[Node]:
        """A gyökértől a megadott számú csomópontig vezető út (Node-lista)."""
        if number not in self._index:
            raise KeyError(f"Ismeretlen csomópont: {number!r}. "
                           f"Elérhető: {sorted(self._index)}")
        return self._index[number]

    def top_level(self) -> list[Node]:
        return self.chapters


# ---------------------------------------------------------------------------
# Breadcrumb XML generálás
# ---------------------------------------------------------------------------

def _make_para(number, title, level, is_current=False):
    """
    Egy breadcrumb-sor (<a:p>) XML-ként.
    level: 0/1/2 -> behúzás (tab). is_current: kiemelés.
    """
    tabs = "\t" * level
    bold = ' b="1"' if is_current else ' b="0"'
    # Aktuális csomópont halvány narancs kiemeléssel (highlight)
    hl = (f'<a:highlight><a:srgbClr val="{C_CURRENT_BG}"/></a:highlight>'
          if is_current else '')

    def run(text, color):
        return (f'<a:r><a:rPr lang="hu-HU" sz="1100"{bold} dirty="0">'
                f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                f'{hl}'
                f'<a:latin typeface="Calibri"/></a:rPr>'
                f'<a:t>{_esc(text)}</a:t></a:r>')

    num_run = run(f"{tabs}{number}. ", C_ORANGE)
    title_run = run(title, C_DARK)
    return (f'<a:p xmlns:a="{NS_A}">'
            f'<a:pPr algn="l" defTabSz="180000">'
            f'<a:spcBef><a:spcPts val="300"/></a:spcBef>'
            f'<a:buNone/>'
            f'</a:pPr>'
            f'{num_run}{title_run}</a:p>')


def build_breadcrumb_xml(tree: ChapterTree, current_number: str) -> list[str]:
    """
    Az aktuális csomóponthoz tartozó breadcrumb paragrafusok (XML stringek).
    A gyökértől az aktuális csomópontig vezető utat rendereli,
    az utolsó (aktuális) elemet kiemelve.
    """
    path = tree.path_to(current_number)
    paras = []
    for i, node in enumerate(path):
        is_current = (i == len(path) - 1)
        paras.append(_make_para(node.number, node.title, node.level, is_current))
    return paras


def build_toc_mindmap_xml(tree: ChapterTree) -> list[str]:
    """A TOC diához: a felső szintű fejezetek listája (nincs kiemelés)."""
    return [_make_para(ch.number, ch.title, 0, False) for ch in tree.top_level()]


def _esc(text: str) -> str:
    return (text.replace("&", "&amp;").replace("<", "&lt;")
                .replace(">", "&gt;").replace('"', "&quot;"))


# ---------------------------------------------------------------------------
# Shape-kezelés
# ---------------------------------------------------------------------------

def _get_shape(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def _find_mindmap_shape(slide):
    """
    Megkeresi a mindmap sidebar shape-et. A mindmap template-ben ez a
    legjobboldalibb 'content_body' (vagy 'mindmap_body' az új .potx-ben).
    """
    candidates = []
    for shape in slide.shapes:
        if shape.name in ("mindmap_body", "content_body") and shape.has_text_frame:
            candidates.append(shape)
    if not candidates:
        return None
    # A legjobboldalibb (legnagyobb left) a mindmap sidebar
    return max(candidates, key=lambda s: s.left or 0)


def _set_mindmap_paragraphs(shape, paras_xml: list[str]):
    """A shape txBody-jának paragrafusait lecseréli a breadcrumb-ra."""
    txBody = shape.text_frame._txBody
    # Töröljük a meglévő <a:p> elemeket
    for p in txBody.findall(f"{{{NS_A}}}p"):
        txBody.remove(p)
    # Beillesztjük az újakat
    for px in paras_xml:
        txBody.append(etree.fromstring(px))


# ---------------------------------------------------------------------------
# Fő API
# ---------------------------------------------------------------------------

def fill_mindmap(template_path: str, output_path: str,
                 tree: ChapterTree, positions: dict[int, str],
                 toc_slides: list[int] | None = None):
    """
    A mindmap template diáit feltölti az auto-generált breadcrumb-bal.

    template_path : a mindmap .pptx template
    output_path   : kimeneti .pptx
    tree          : ChapterTree — a fejezet-hierarchia
    positions     : {slide_index (0-alapú): csomópont-szám} — melyik dia
                    hol tart a fában. Pl. {3: "1.1.1"}.
    toc_slides    : azon diák indexei, ahol a TELJES fejezetlista jelenjen
                    meg a sidebarban (nem breadcrumb). Pl. [1] a TOC dia.
    """
    toc_slides = toc_slides or []
    prs = Presentation(template_path)

    filled = 0
    for idx, slide in enumerate(prs.slides):
        shape = _find_mindmap_shape(slide)
        if shape is None:
            continue

        if idx in toc_slides:
            paras = build_toc_mindmap_xml(tree)
        elif idx in positions:
            paras = build_breadcrumb_xml(tree, positions[idx])
        else:
            continue

        _set_mindmap_paragraphs(shape, paras)
        filled += 1

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"Mindmap feltöltve {filled} dián. Mentve: {output_path}")
    return prs


# ---------------------------------------------------------------------------
# CLI demó
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import sys

    template = sys.argv[1] if len(sys.argv) > 1 else \
        "templates/due_prenetation_template_mindmap.pptx"
    output = sys.argv[2] if len(sys.argv) > 2 else \
        "output/mindmap_demo.pptx"

    # Példa fejezet-fa
    tree = ChapterTree([
        Chapter("1", "Első fejezet", [
            Section("1.1", "Első szakasz", [
                "1.1.1 Első alszakasz",
                "1.1.2 Második alszakasz",
            ]),
            Section("1.2", "Második szakasz"),
        ]),
        Chapter("2", "Második fejezet", [
            Section("2.1", "Második szakasz"),
        ]),
        Chapter("3", "Harmadik fejezet"),
        Chapter("4", "Negyedik fejezet"),
    ])

    # Diánkénti pozíció (0-alapú slide-index -> csomópont)
    # A mindmap template: 0=cím, 1=TOC, 2=szakaszfejléc, 3=tartalom, ...
    positions = {
        2: "1",        # szakaszfejléc -> 1. fejezet
        3: "1.1.1",    # tartalom dia -> 1.1.1 alszakasz
        4: "1.1.1",    # kép+szöveg -> 1.1.1
        5: "1.1.1",    # ábra -> 1.1.1
        6: "1.1.1",    # táblázat -> 1.1.1
    }
    toc_slides = [1]   # a TOC dián a teljes fejezetlista

    fill_mindmap(template, output, tree, positions, toc_slides)
    print("Kész.")
