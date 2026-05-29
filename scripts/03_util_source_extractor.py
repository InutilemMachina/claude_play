"""
03_util_source_extractor.py -- Nem-PDF források determinisztikus extraktora.

A MinerU csak PDF-et kezel. Ez a script a többi forrástípust dolgozza fel
MinerU-kompatibilis kimeneti struktúrába:
    2_clean_inputs/<stem>/auto/<stem>.md       (kinyert szöveg, Markdown)
    2_clean_inputs/<stem>/auto/images/*.{png,jpg}  (beágyazott képek, ha vannak)

Támogatott típusok:
    .pptx  -- python-pptx (dia-szöveg + beágyazott képek)
    .html  -- beautifulsoup4 (törzsszöveg, script/style/nav eltávolítva)
    .docx  -- python-docx (ha telepítve; egyébként figyelmeztet és kihagy)

A PDF fájlokat kihagyja (azokat a MinerU dolgozza fel).

Usage:
    python scripts/03_util_source_extractor.py --week-dir <path/to/N_het>
    python scripts/03_util_source_extractor.py --week-dir <path> --types pptx html
"""

import argparse
import sys
from pathlib import Path

try:
    from _encoding_fix import fix_stdout as _fix_stdout
    _fix_stdout()
except ImportError:
    pass


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def extract_pptx(src: Path, out_dir: Path) -> tuple[str, int]:
    """Extract slide text + images from a .pptx. Returns (markdown, n_images)."""
    from pptx import Presentation
    from pptx.util import Emu  # noqa: F401

    prs = Presentation(str(src))
    img_dir = out_dir / "images"
    lines = [f"# {src.stem}", ""]
    n_images = 0

    for idx, slide in enumerate(prs.slides, 1):
        lines.append(f"## Dia {idx}")
        lines.append("")
        for shape in slide.shapes:
            # Text
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        # Heuristic: short bold-ish lines as bullets
                        lines.append(f"- {txt}")
            # Tables
            if shape.has_table:
                tbl = shape.table
                rows = list(tbl.rows)
                if rows:
                    header = [c.text.strip() for c in rows[0].cells]
                    lines.append("")
                    lines.append("| " + " | ".join(header) + " |")
                    lines.append("| " + " | ".join("---" for _ in header) + " |")
                    for r in rows[1:]:
                        lines.append("| " + " | ".join(c.text.strip() for c in r.cells) + " |")
                    lines.append("")
            # Images
            if shape.shape_type == 13:  # PICTURE
                try:
                    image = shape.image
                    ext = image.ext
                    img_dir.mkdir(parents=True, exist_ok=True)
                    n_images += 1
                    img_name = f"slide{idx}_img{n_images}.{ext}"
                    (img_dir / img_name).write_bytes(image.blob)
                    lines.append("")
                    lines.append(f"![Dia {idx} kép]({img_name})")
                    lines.append("")
                except Exception:
                    pass
        lines.append("")

    return "\n".join(lines), n_images


# ---------------------------------------------------------------------------
# HTML
# ---------------------------------------------------------------------------

def extract_html(src: Path, out_dir: Path) -> tuple[str, int]:
    """Extract main text from a local .html file. Returns (markdown, n_images)."""
    from bs4 import BeautifulSoup

    raw = src.read_bytes().decode("utf-8", errors="replace")
    soup = BeautifulSoup(raw, "html.parser")

    # Strip non-content elements
    for tag in soup(["script", "style", "nav", "footer", "header", "aside",
                     "noscript", "form", "iframe"]):
        tag.decompose()

    title = (soup.title.string.strip() if soup.title and soup.title.string
             else src.stem)
    lines = [f"# {title}", ""]

    # Walk headings + paragraphs + list items in document order
    body = soup.body or soup
    for el in body.find_all(["h1", "h2", "h3", "h4", "p", "li"]):
        text = el.get_text(" ", strip=True)
        if not text or len(text) < 2:
            continue
        name = el.name
        if name in ("h1", "h2"):
            lines.append(f"## {text}")
        elif name in ("h3", "h4"):
            lines.append(f"### {text}")
        elif name == "li":
            lines.append(f"- {text}")
        else:
            lines.append(text)
        lines.append("")

    # Note: local HTML rarely has downloadable images (often data: URIs or
    # external) -- image extraction skipped for HTML (use C.3/MinerU for figures).
    return "\n".join(lines), 0


# ---------------------------------------------------------------------------
# DOCX (optional -- requires python-docx)
# ---------------------------------------------------------------------------

def extract_docx(src: Path, out_dir: Path) -> tuple[str, int]:
    """Extract text from .docx. Returns (markdown, n_images)."""
    try:
        import docx
    except ImportError:
        raise RuntimeError(
            "python-docx nincs telepítve. Telepítés: pip install python-docx "
            "(vagy konvertáld a DOCX-et PDF-re és használd a MinerU-t)."
        )

    document = docx.Document(str(src))
    lines = [f"# {src.stem}", ""]
    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower()
        if "heading 1" in style or "title" in style:
            lines.append(f"## {text}")
        elif "heading" in style:
            lines.append(f"### {text}")
        elif style.startswith("list"):
            lines.append(f"- {text}")
        else:
            lines.append(text)
        lines.append("")
    return "\n".join(lines), 0


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

EXTRACTORS = {
    "pptx": extract_pptx,
    "html": extract_html,
    "docx": extract_docx,
}


def main():
    parser = argparse.ArgumentParser(description="Nem-PDF források extraktora (MinerU-kompatibilis kimenet)")
    parser.add_argument("--week-dir", required=True, type=Path)
    parser.add_argument("--types", nargs="+", default=["pptx", "html", "docx"],
                        choices=["pptx", "html", "docx"],
                        help="Feldolgozandó forrástípusok (default: mind)")
    parser.add_argument("--dry-run", action="store_true")
    args = parser.parse_args()

    week_dir = args.week_dir.resolve()
    raw_in = week_dir / "1_raw_inputs"
    clean_in = week_dir / "2_clean_inputs"
    if not raw_in.is_dir():
        sys.exit(f"HIBA: nem található {raw_in}")

    n_ok = 0
    n_skip = 0
    for src in sorted(raw_in.iterdir()):
        if not src.is_file():
            continue
        ext = src.suffix.lower().lstrip(".")
        if ext == "pdf":
            continue  # MinerU dolgozza fel
        if ext not in args.types:
            continue
        extractor = EXTRACTORS.get(ext)
        if not extractor:
            continue

        out_dir = clean_in / src.stem / "auto"
        try:
            md, n_imgs = extractor(src, out_dir)
        except RuntimeError as e:
            print(f"  SKIP  {src.name}: {e}", file=sys.stderr)
            n_skip += 1
            continue
        except Exception as e:
            print(f"  HIBA  {src.name}: {type(e).__name__}: {e}", file=sys.stderr)
            n_skip += 1
            continue

        if args.dry_run:
            print(f"  [DRY] {src.name} -> {len(md)} kar, {n_imgs} kép")
        else:
            out_dir.mkdir(parents=True, exist_ok=True)
            (out_dir / f"{src.stem}.md").write_text(md, encoding="utf-8")
            print(f"  OK    {src.name} -> 2_clean_inputs/{src.stem}/auto/{src.stem}.md "
                  f"({len(md)} kar, {n_imgs} kép)")
        n_ok += 1

    print(f"Kész: {n_ok} forrás feldolgozva, {n_skip} kihagyva.")


if __name__ == "__main__":
    main()
