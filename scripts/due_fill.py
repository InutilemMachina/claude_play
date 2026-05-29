"""
due_fill.py  --  DUE prezentáció template kitöltő
==============================================

Használat:
    from scripts.due_fill import DUEPresentation

    prs = DUEPresentation("templates/due_refactored.pptx")
    prs.set_title("Bevezetés a biofizikába", "Fizioterápiás BSc • II. évfolyam")
    prs.set_global_footer("Dr. Hári László", "2026.09.01.")

    prs.add_section("01", "Fizikai alapok", "A mechanika és termodinamika rövid összefoglalója.")

    idx = prs.add_content_slide("Mechanika alapjai", level="h1")
    prs.set_bullets(idx, [
        ("Sebesség és gyorsulás", "h1"),
        ("Newton törvények", "h2"),
        ("Munka és energia", "h2"),
    ])

    idx = prs.add_table_slide("Mértékegységek", "1. táblázat: SI mértékegységek")
    prs.set_table(idx, ["Mennyiség","Jel","Egység"], [
        ["Tömeg","m","kg"],
        ["Erő","F","N"],
    ])

    prs.save("output/biofizika.pptx")

Megjegyzés:
    A template-ben lévő placeholder slide-ok automatikusan törlődnek a save() előtt.
    A section-ok és content slide-ok page-számozása mentéskor frissül.
"""

from __future__ import annotations
import copy
import os
from pathlib import Path
from typing import Literal

try:
    from pptx import Presentation
    from pptx.util import Pt, Emu, Inches
    from pptx.dml.color import RGBColor
    from lxml import etree
except ImportError:
    raise ImportError("pip install python-pptx lxml")


# ---------------------------------------------------------------------------
# Típusok
# ---------------------------------------------------------------------------

Level = Literal["h1", "h2", "h3"]

# Dia-típus → template-index (0-alapú) a due_refactored.pptx-ben
TEMPLATE_IDX = {
    "title":      0,   # CÍM DIA
    "toc":        1,   # TARTALOM DIA (tartalomjegyzék)
    "section":    2,   # SZAKASZFEJLÉC DIA
    "h1":         3,   # TARTALOM szint-1 (fejezet cím)
    "h2":         4,   # TARTALOM szint-2 (szakasz)
    "h3":         5,   # TARTALOM szint-3 (alszakasz)
    "image_text": 6,   # KÉP + SZÖVEG
    "figure":     7,   # ÁBRA
    "table":      8,   # TÁBLÁZAT
    "refs":       9,   # IRODALOMJEGYZÉK
    "blank":      10,  # ÜRES TARTALOM
    "changelog":  11,  # VÁLTOZÁSJEGYZÉK
}

# Orange szám szín és dark body szín (Calibri)
C_ORANGE  = "ED7D31"
C_DARK    = "212121"
C_NAVY    = "0D1B3E"
C_GRAY    = "7A8A9E"
C_WHITE   = "FFFFFF"

# Bekezdés szóköz és egyéb formázás konstansok (EMU / hundredths of a point)
BULLET_SPC_BEF  = 400   # pts/100
BULLET_SPC_AFT  = 100
BULLET_TAB_SZ   = 180000


# ---------------------------------------------------------------------------
# Segédfüggvények
# ---------------------------------------------------------------------------

def _get_shape(slide, name: str):
    """Névvel keres shape-t a dián; None ha nem találja."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def _set_single_run_text(slide, shape_name: str, text: str) -> bool:
    """
    Lecseréli a shape szövegét `text`-re, az első run formázását megőrizve.
    Minden többi runt és paragrafust eltávolít.
    """
    shape = _get_shape(slide, shape_name)
    if not shape or not shape.has_text_frame:
        return False
    tf = shape.text_frame
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"

    # Megőrizzük az első paramétereket (pPr) és az első run formátumát (rPr)
    first_ppr = None
    first_rpr = None
    for para in tf.paragraphs:
        ppr = para._p.find(f"{{{ns}}}pPr")
        if ppr is not None:
            first_ppr = copy.deepcopy(ppr)
        for run in para.runs:
            rpr = run._r.find(f"{{{ns}}}rPr")
            if rpr is not None:
                first_rpr = copy.deepcopy(rpr)
            break
        break

    # Töröljük az összes paragrafust
    _clear_text_frame(shape)

    # Új paragrafus az eredeti formázással
    p = etree.SubElement(tf._txBody, f"{{{ns}}}p")
    if first_ppr is not None:
        p.insert(0, copy.deepcopy(first_ppr))
    r = etree.SubElement(p, f"{{{ns}}}r")
    if first_rpr is not None:
        r.insert(0, copy.deepcopy(first_rpr))
    t_elem = etree.SubElement(r, f"{{{ns}}}t")
    t_elem.text = text
    return True


def _clear_text_frame(shape):
    """Az összes paragrafust és runt törli, de a bodyPr megmarad."""
    tf = shape.text_frame
    spTree = tf._txBody
    # Remove all <a:p> elements
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for p in spTree.findall(f"{{{ns}}}p"):
        spTree.remove(p)


def _make_run_xml(text: str, sz: int, bold: bool, color_hex: str,
                  lang: str = "hu-HU", typeface: str = "Calibri") -> str:
    """Egy <a:r> XML stringet épít fel."""
    b_attr = ' b="1"' if bold else ' b="0"'
    return (
        f'<a:r xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:rPr lang="{lang}" sz="{sz}"{b_attr} dirty="0">'
        f'<a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill>'
        f'<a:latin typeface="{typeface}" panose="020B0004020202020204" pitchFamily="34" charset="0"/>'
        f'</a:rPr>'
        f'<a:t>{_esc(text)}</a:t>'
        f'</a:r>'
    )


def _make_para_xml(runs_xml: str, align: str = "l",
                   spc_bef: int = BULLET_SPC_BEF,
                   spc_aft: int = BULLET_SPC_AFT,
                   tab_sz: int = BULLET_TAB_SZ) -> str:
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    return (
        f'<a:p xmlns:a="{ns}">'
        f'<a:pPr algn="{align}" defTabSz="{tab_sz}">'
        f'<a:spcBef><a:spcPts val="{spc_bef}"/></a:spcBef>'
        f'<a:spcAft><a:spcPts val="{spc_aft}"/></a:spcAft>'
        f'</a:pPr>'
        f'{runs_xml}'
        f'</a:p>'
    )


def _esc(text: str) -> str:
    """XML entity escaping."""
    return (text.replace("&", "&amp;")
                .replace("<", "&lt;")
                .replace(">", "&gt;")
                .replace('"', "&quot;"))


def _append_xml_to_txbody(shape, xml_str: str):
    """Egy XML string-et parsed elem-ként hozzáfűz a txBody-hoz."""
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    elem = etree.fromstring(xml_str)
    shape.text_frame._txBody.append(elem)


def _clone_slide(prs: Presentation, source_idx: int):
    """
    Klónoz egy template slide-ot és a prezentáció végéhez adja.
    Visszatérési érték: az új Slide objektum.
    """
    src = prs.slides[source_idx]
    slide_layout = src.slide_layout

    new_slide = prs.slides.add_slide(slide_layout)

    # Copy shape tree (skip first 2: nvGrpSpPr, grpSpPr)
    sp_src = src.shapes._spTree
    sp_dst = new_slide.shapes._spTree
    for elem in list(sp_dst)[2:]:
        sp_dst.remove(elem)
    for elem in list(sp_src)[2:]:
        sp_dst.append(copy.deepcopy(elem))

    # Copy image relationships preserving the original rId (so r:embed refs stay valid)
    for rId, rel in src.part._rels._rels.items():
        if "image" in rel.reltype and rId not in new_slide.part._rels._rels:
            new_slide.part._rels._rels[rId] = rel

    return new_slide


# ---------------------------------------------------------------------------
# Fő osztály
# ---------------------------------------------------------------------------

class DUEPresentation:
    """DUE template-alapú prezentáció builder."""

    def __init__(self, template_path: str):
        self.prs = Presentation(template_path)
        self._author = ""
        self._date = ""
        self._output_slides: list[dict] = []  # {slide, type, page_visible}
        self._page_counter = 1  # oldalsorszám (1-ről indul a title utáni dia)

    # ------------------------------------------------------------------
    # Globális beállítások
    # ------------------------------------------------------------------

    def set_global_footer(self, author: str, date: str):
        """Beállítja az összes dia footer_author és footer_date értékét."""
        self._author = author
        self._date = date

    # ------------------------------------------------------------------
    # Cím dia
    # ------------------------------------------------------------------

    def set_title(self, main_title: str, subtitle: str = ""):
        """Kitölti a template 1. diáját (cím dia); nem adja hozzá az output-hoz."""
        slide = self.prs.slides[TEMPLATE_IDX["title"]]
        _set_single_run_text(slide, "slide_main_title", main_title)
        if subtitle:
            _set_single_run_text(slide, "slide_subtitle", subtitle)
        if self._author:
            _set_single_run_text(slide, "footer_author", self._author)
        if self._date:
            _set_single_run_text(slide, "footer_date", self._date)
        self._output_slides.insert(0, {"slide": slide, "type": "title", "page_visible": False})

    # ------------------------------------------------------------------
    # Tartalomjegyzék dia
    # ------------------------------------------------------------------

    def add_toc(self, title: str, entries: list[tuple[str, Level]]) -> int:
        """
        Tartalomjegyzék diát ad hozzá.
        entries: [(szöveg, szint)] pl. [("1. Fejezet", "h1"), ("1.1. Szakasz", "h2")]
        Visszatérési érték: az új dia indexe az output listában.
        """
        slide = _clone_slide(self.prs, TEMPLATE_IDX["toc"])
        self._fill_header_footer(slide, title)
        self._fill_hierarchical_body(slide, "content_body", entries)
        return self._register(slide, "toc")

    # ------------------------------------------------------------------
    # Szakaszfejléc dia
    # ------------------------------------------------------------------

    def add_section(self, number: str, title: str, description: str = "") -> int:
        """Szakaszfejléc diát ad hozzá (nem kap oldalszámot)."""
        slide = _clone_slide(self.prs, TEMPLATE_IDX["section"])
        _set_single_run_text(slide, "section_number", number)
        _set_single_run_text(slide, "section_title", title)
        if description:
            _set_single_run_text(slide, "section_description", description)
        return self._register(slide, "section", page_visible=False)

    # ------------------------------------------------------------------
    # Tartalom dia (bullet lista)
    # ------------------------------------------------------------------

    def add_content_slide(self, title: str,
                          bullets: list[tuple[str, Level]] | None = None,
                          level: Level = "h1") -> int:
        """
        Bullet lista diát ad hozzá.
        bullets: [(szöveg, szint)] — ha None, üres content_body marad.
        level: meghatározza a template forrás-diát ("h1"/"h2"/"h3").
        """
        slide = _clone_slide(self.prs, TEMPLATE_IDX[level])
        self._fill_header_footer(slide, title)
        if bullets:
            self._fill_bullet_body(slide, "content_body", bullets)
        return self._register(slide, "content")

    # ------------------------------------------------------------------
    # Kép + szöveg dia
    # ------------------------------------------------------------------

    def add_image_text_slide(self, title: str,
                             bullets: list[tuple[str, Level]],
                             image_path: str,
                             caption: str = "") -> int:
        """Kép + szöveg diát ad hozzá."""
        slide = _clone_slide(self.prs, TEMPLATE_IDX["image_text"])
        self._fill_header_footer(slide, title)
        if bullets:
            self._fill_bullet_body(slide, "col_left_body", bullets)
        if caption:
            _set_single_run_text(slide, "figure_caption", caption)
        if image_path and os.path.exists(image_path):
            self._replace_image_placeholder(slide, "col_right_image", image_path)
        return self._register(slide, "image_text")

    # ------------------------------------------------------------------
    # Ábra dia
    # ------------------------------------------------------------------

    def add_figure_slide(self, title: str, image_path: str, caption: str = "") -> int:
        """Teljes-szélességű ábra diát ad hozzá."""
        slide = _clone_slide(self.prs, TEMPLATE_IDX["figure"])
        self._fill_header_footer(slide, title)
        if caption:
            _set_single_run_text(slide, "figure_caption", caption)
        if image_path and os.path.exists(image_path):
            self._replace_image_placeholder(slide, "figure_image", image_path)
        return self._register(slide, "figure")

    # ------------------------------------------------------------------
    # Táblázat dia
    # ------------------------------------------------------------------

    def add_table_slide(self, title: str, table_title: str,
                        headers: list[str],
                        rows: list[list[str]]) -> int:
        """Táblázat diát ad hozzá és feltölti az adatokkal."""
        slide = _clone_slide(self.prs, TEMPLATE_IDX["table"])
        self._fill_header_footer(slide, title)
        _set_single_run_text(slide, "table_title", table_title)
        self._fill_table(slide, "Table 10", headers, rows)
        return self._register(slide, "table")

    # ------------------------------------------------------------------
    # Irodalomjegyzék dia
    # ------------------------------------------------------------------

    def add_refs_slide(self, refs: list[str]) -> int:
        """Irodalomjegyzék diát ad hozzá."""
        slide = _clone_slide(self.prs, TEMPLATE_IDX["refs"])
        self._fill_header_footer(slide, "Irodalomjegyzék")
        entries = [(ref, "h1") for ref in refs]
        self._fill_hierarchical_body(slide, "content_body", entries)
        return self._register(slide, "refs")

    # ------------------------------------------------------------------
    # Változásjegyzék dia
    # ------------------------------------------------------------------

    def add_changelog(self, entries: list[tuple[str, str, str, str]]) -> int:
        """
        Változásjegyzék diát ad hozzá.
        entries: [(verzió, dátum, szerző, leírás)]
        """
        slide = _clone_slide(self.prs, TEMPLATE_IDX["changelog"])
        self._fill_header_footer(slide, "Változásjegyzék")
        self._fill_changelog_table(slide, entries)
        return self._register(slide, "changelog")

    # ------------------------------------------------------------------
    # Mentés
    # ------------------------------------------------------------------

    def save(self, output_path: str):
        """
        Véglegesíti a prezentációt és elmenti.
        1. Oldalszámokat frissíti az output slide-okon.
        2. A template placeholder slide-okat eltávolítja.
        3. Elmenti a PPTX-et.
        """
        self._update_page_numbers()
        self._apply_global_footer()
        self._remove_template_slides()
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(output_path)
        print(f"Mentve: {output_path}")

    # ------------------------------------------------------------------
    # Belső segédmetódusok
    # ------------------------------------------------------------------

    def _register(self, slide, slide_type: str, page_visible: bool = True) -> int:
        """Regisztrálja az output slide-ot és visszaadja a listabeli indexet."""
        self._output_slides.append({
            "slide": slide,
            "type": slide_type,
            "page_visible": page_visible,
        })
        return len(self._output_slides) - 1

    def _fill_header_footer(self, slide, title: str):
        """slide_title, footer_date, footer_author kitöltése."""
        _set_single_run_text(slide, "slide_title", title)
        if self._date:
            _set_single_run_text(slide, "footer_date", self._date)
        if self._author:
            _set_single_run_text(slide, "footer_author", self._author)

    def _fill_bullet_body(self, slide, shape_name: str,
                          bullets: list[tuple[str, Level]]):
        """
        Feltölti a content_body text box-ot bullet-listával.
        bullets: [(szöveg, szint)] ahol szint in {"h1","h2","h3"}
        """
        shape = _get_shape(slide, shape_name)
        if not shape:
            return
        _clear_text_frame(shape)

        for text, level in bullets:
            if level == "h1":
                prefix = "▶ "
                sz = 1500
                tab_prefix = ""
            elif level == "h2":
                prefix = "–  "
                sz = 1400
                tab_prefix = "\t"
            else:
                prefix = "·  "
                sz = 1300
                tab_prefix = "\t\t"

            run_xml = _make_run_xml(tab_prefix + prefix + text, sz, False, C_DARK)
            para_xml = _make_para_xml(run_xml)
            _append_xml_to_txbody(shape, para_xml)

    def _fill_hierarchical_body(self, slide, shape_name: str,
                                entries: list[tuple[str, Level]]):
        """
        Tartalomjegyzék / irodalomjegyzék jellegű szöveg: narancs prefix + sötét szöveg.
        entries: [(szöveg, szint)]
        """
        shape = _get_shape(slide, shape_name)
        if not shape:
            return
        _clear_text_frame(shape)

        for text, level in entries:
            if level == "h1":
                tab = ""
                sz = 1600
            elif level == "h2":
                tab = "\t"
                sz = 1600
            else:
                tab = "\t\t"
                sz = 1600

            # Narancs prefix-run + sötét szöveg-run
            # A szöveg első "szó" (pl. "1.", "1.1.") narancs lesz, a többi sötét
            parts = text.split(" ", 1)
            prefix_txt = parts[0] + " "
            body_txt = parts[1] if len(parts) > 1 else ""

            run1 = _make_run_xml(tab + prefix_txt, sz, False, C_ORANGE, typeface="Aptos")
            run2 = _make_run_xml(body_txt, sz, False, C_DARK, typeface="Aptos") if body_txt else ""
            para_xml = _make_para_xml(run1 + run2)
            _append_xml_to_txbody(shape, para_xml)

    def _fill_table(self, slide, shape_name: str,
                    headers: list[str], rows: list[list[str]]):
        """A template táblázatot feltölti adatokkal."""
        shape = _get_shape(slide, shape_name)
        if not shape or not shape.has_table:
            return
        table = shape.table
        n_rows, n_cols = len(table.rows), len(table.columns)

        def set_cell(r, c, text):
            if r < n_rows and c < n_cols:
                cell = table.cell(r, c)
                cell.text_frame.paragraphs[0].runs[0].text = text if \
                    cell.text_frame.paragraphs[0].runs else ""
                # Ha nincs run, adunk egyet
                if not cell.text_frame.paragraphs[0].runs:
                    run = cell.text_frame.paragraphs[0].add_run()
                    run.text = text

        for c, hdr in enumerate(headers[:n_cols]):
            set_cell(0, c, hdr)
        for r, row in enumerate(rows[:n_rows - 1]):
            for c, val in enumerate(row[:n_cols]):
                set_cell(r + 1, c, val)

    def _fill_changelog_table(self, slide, entries: list[tuple[str, str, str, str]]):
        """A változásjegyzék táblázatot feltölti."""
        shape = _get_shape(slide, "content_body")
        if not shape:
            return
        _clear_text_frame(shape)
        # Fejléc sor
        hdr_text = "Verzió\tDátum\tSzerző\tLeírás"
        run_hdr = _make_run_xml(hdr_text, 1300, True, C_NAVY)
        _append_xml_to_txbody(shape, _make_para_xml(run_hdr))
        # Adatsorok
        for ver, date, author, desc in entries:
            row_text = f"{ver}\t{date}\t{author}\t{desc}"
            run_row = _make_run_xml(row_text, 1300, False, C_DARK)
            _append_xml_to_txbody(shape, _make_para_xml(run_row))

    def _replace_image_placeholder(self, slide, shape_name: str, image_path: str):
        """Lecseréli az image placeholder shape-et egy valódi képpel."""
        shape = _get_shape(slide, shape_name)
        if not shape:
            return
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        sp = shape._element
        sp.getparent().remove(sp)
        slide.shapes.add_picture(image_path, left, top, width, height)

    def _update_page_numbers(self):
        """Az output slide-ok footer_page mezőit frissíti."""
        page = 1
        for entry in self._output_slides:
            if entry["type"] == "title":
                continue
            if entry["page_visible"]:
                _set_single_run_text(entry["slide"], "footer_page", str(page))
                page += 1
            else:
                # Szakaszfejléc: oldalszám nem látható, de a counter nem ugrik
                pass

    def _apply_global_footer(self):
        """A globális footer értékeket alkalmazza az összes output slide-ra."""
        for entry in self._output_slides:
            slide = entry["slide"]
            if self._author:
                _set_single_run_text(slide, "footer_author", self._author)
            if self._date:
                _set_single_run_text(slide, "footer_date", self._date)

    def _remove_template_slides(self):
        """
        Eltávolítja azokat a template slide-okat, amelyek nem kerültek
        az output_slides listába.
        """
        ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        output_parts = {e["slide"].part for e in self._output_slides}
        sldIdLst = self.prs.slides._sldIdLst
        prs_part = self.prs.part

        for sldId in list(sldIdLst):
            rid = sldId.get(f"{{{ns_r}}}id")
            slide_part = prs_part.related_part(rid)
            if slide_part not in output_parts:
                sldIdLst.remove(sldId)


# ---------------------------------------------------------------------------
# CLI teszt / példa
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import sys

    template = sys.argv[1] if len(sys.argv) > 1 else "templates/due_refactored.pptx"
    output   = sys.argv[2] if len(sys.argv) > 2 else "output/due_test_output.pptx"

    prs = DUEPresentation(template)
    prs.set_global_footer("Dr. Hári László", "2026.09.01.")

    prs.set_title(
        "Bevezetés a biofizikába",
        "Fizioterápiás BSc • Biofizika • DE GYKK"
    )

    prs.add_toc("Tartalom", [
        ("1. Fizikai alapok", "h1"),
        ("1.1. Mechanika", "h2"),
        ("1.2. Termodinamika", "h2"),
        ("2. Hullámtan", "h1"),
        ("2.1. Hullámok típusai", "h2"),
    ])

    prs.add_section("01", "Fizikai alapok",
                    "A mechanika és termodinamika rövid összefoglalója.")

    prs.add_content_slide("Mechanika alapjai", [
        ("Newton I. törvénye: tehetetlenség", "h1"),
        ("Newton II. törvénye: F = ma", "h1"),
        ("Newton III. törvénye: hatás-ellenhatás", "h1"),
    ], level="h1")

    prs.add_content_slide("Kinematika", [
        ("Egyenes vonalú egyenletes mozgás", "h2"),
        ("Egyenletesen változó mozgás", "h2"),
        ("Körmozgás és szögsebesség", "h2"),
    ], level="h2")

    prs.add_table_slide(
        title="Fizikai mennyiségek",
        table_title="1. táblázat: Alapvető SI mennyiségek",
        headers=["Mennyiség", "Jel", "Mértékegység"],
        rows=[
            ["Tömeg", "m", "kg"],
            ["Erő", "F", "N (Newton)"],
            ["Energia", "E", "J (Joule)"],
            ["Nyomás", "p", "Pa (Pascal)"],
        ]
    )

    prs.add_refs_slide([
        "[1] Atkins, P. (2014). Physical Chemistry. Oxford University Press.",
        "[2] Halliday, D. et al. (2013). Fundamentals of Physics. Wiley.",
        "[3] Nelson, P. (2008). Biological Physics. Freeman.",
    ])

    prs.save(output)
    print(f"Kész: {output}")
