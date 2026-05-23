"""
pptx_gyarto.py -- Marp Markdown → PPTX converter
Converts a Marp-formatted Markdown file to a PowerPoint presentation.

Usage:
    python scripts/pptx_gyarto.py <marp_md_file> [--template <template.pptx>] [--output <out.pptx>]

If --template is not provided (or file missing), uses a built-in default style.

NOTE (pipeline): du_template.pptx is the intended template. If missing,
this script generates a functionally equivalent default-styled deck.
"""

import argparse
import re
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    sys.exit("python-pptx not installed. Run: pip install python-pptx --break-system-packages")


# ---------------------------------------------------------------------------
# Slide dimensions: 16:9 widescreen
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Colours (neutral academic theme)
C_TITLE_BG  = RGBColor(0x1F, 0x49, 0x7D)   # dark blue
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK     = RGBColor(0x1A, 0x1A, 0x1A)
C_ACCENT    = RGBColor(0x2E, 0x75, 0xB6)   # medium blue
C_LIGHT_BG  = RGBColor(0xF2, 0xF2, 0xF2)   # light grey
C_MSC_BG    = RGBColor(0xFF, 0xF0, 0xD0)   # warm amber for MSc slides


def parse_marp(md_text: str) -> list[dict]:
    """
    Parse Marp Markdown into a list of slide dicts.
    Each dict: {title, subtitle, body, is_msc, is_cover}
    """
    # Strip YAML frontmatter
    md_text = re.sub(r'^---\n.*?---\n', '', md_text, count=1, flags=re.DOTALL)

    # Split on --- slide separators
    raw_slides = re.split(r'\n---\n', md_text)

    slides = []
    for raw in raw_slides:
        raw = raw.strip()
        if not raw:
            continue

        is_msc = bool(re.search(r'<!-- MSc -->', raw))
        # Strip MSc comments
        raw = re.sub(r'<!-- /?MSc -->', '', raw).strip()

        lines = raw.split('\n')
        title = ''
        subtitle = ''
        body_lines = []
        skip_next = False

        for i, line in enumerate(lines):
            if skip_next:
                skip_next = False
                continue
            if re.match(r'^# ', line):
                title = line[2:].strip()
                # Next line may be subtitle (###)
                if i + 1 < len(lines) and re.match(r'^### ', lines[i+1]):
                    subtitle = lines[i+1][4:].strip()
                    skip_next = True
            elif re.match(r'^## ', line):
                title = line[3:].strip()
            elif re.match(r'^### ', line):
                subtitle = line[4:].strip()
            else:
                body_lines.append(line)

        body = '\n'.join(body_lines).strip()
        slides.append({
            'title': title,
            'subtitle': subtitle,
            'body': body,
            'is_msc': is_msc,
            'is_cover': bool(not subtitle and re.match(r'^# ', raw.split('\n')[0])) if lines else False,
        })
    return slides


def add_textbox(slide, left, top, width, height, text, font_size=Pt(18),
                bold=False, color=None, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def fill_shape(shape, color: RGBColor):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def build_presentation(slides_data: list[dict], template_path: str | None) -> Presentation:
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
        print(f"Template betöltve: {template_path}")
    else:
        prs = Presentation()
        prs.slide_width  = SLIDE_W
        prs.slide_height = SLIDE_H
        if template_path:
            print(f"⚠️  Helyettesítés: {template_path!r} nem található -- "
                  "du_template.pptx hiányában default stílust használok. "
                  "Dokumentáció: kepek_workflow.md §7 Teszt helyettesítés.")

    blank_layout = prs.slide_layouts[6]  # blank

    for idx, sd in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        w, h = prs.slide_width, prs.slide_height
        m = Inches(0.4)   # margin

        # Background colour
        bg_color = C_MSC_BG if sd['is_msc'] else (C_TITLE_BG if idx == 0 else C_WHITE)
        bg = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            0, 0, w, h
        )
        fill_shape(bg, bg_color)
        bg.line.fill.background()

        text_color = C_WHITE if (idx == 0 and not sd['is_msc']) else C_BLACK

        # Title bar (coloured strip at top for non-cover slides)
        if idx > 0:
            bar = slide.shapes.add_shape(1, 0, 0, w, Inches(1.1))
            fill_shape(bar, C_TITLE_BG)
            bar.line.fill.background()

        # Title
        if sd['title']:
            title_top = Inches(0.15) if idx == 0 else Inches(0.1)
            title_h   = Inches(1.0)
            add_textbox(
                slide, m, title_top, w - 2*m, title_h,
                sd['title'],
                font_size=Pt(32) if idx == 0 else Pt(26),
                bold=True,
                color=C_WHITE,
                align=PP_ALIGN.CENTER if idx == 0 else PP_ALIGN.LEFT,
            )

        # Subtitle (cover slide)
        if sd['subtitle'] and idx == 0:
            add_textbox(
                slide, m, Inches(1.4), w - 2*m, Inches(0.6),
                sd['subtitle'],
                font_size=Pt(20), bold=False,
                color=C_WHITE, align=PP_ALIGN.CENTER,
            )

        # MSc badge
        if sd['is_msc']:
            badge = slide.shapes.add_shape(1, w - Inches(1.8), Inches(0.05), Inches(1.7), Inches(0.45))
            fill_shape(badge, RGBColor(0xD0, 0x70, 0x00))
            badge.line.fill.background()
            add_textbox(slide, w - Inches(1.8), Inches(0.06), Inches(1.7), Inches(0.4),
                        "MSc", font_size=Pt(14), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Body
        if sd['body']:
            body_top = Inches(1.25) if idx > 0 else Inches(2.1)
            body_h   = h - body_top - Inches(0.3)

            # Render body: handle tables, code blocks, bullets
            body_text = sd['body']

            # Clean up Markdown table formatting to plain text
            if '|' in body_text:
                lines = body_text.split('\n')
                clean = []
                for ln in lines:
                    if re.match(r'^\s*\|[-: |]+\|\s*$', ln):
                        continue  # separator row
                    ln = re.sub(r'^\s*\|', '', ln)
                    ln = re.sub(r'\|\s*$', '', ln)
                    parts = [p.strip() for p in ln.split('|')]
                    clean.append('  '.join(parts))
                body_text = '\n'.join(clean)

            # Strip code fences
            body_text = re.sub(r'```[^\n]*\n', '', body_text)
            body_text = re.sub(r'```', '', body_text)

            # Strip HTML comments
            body_text = re.sub(r'<!--.*?-->', '', body_text, flags=re.DOTALL)

            # Simplify LaTeX (keep as-is, pptx won't render it, but text stays)
            body_text = body_text.strip()

            # Adjust font size based on length
            n_chars = len(body_text)
            fsize = Pt(14) if n_chars > 600 else Pt(16) if n_chars > 300 else Pt(18)

            add_textbox(
                slide, m, body_top, w - 2*m, body_h,
                body_text,
                font_size=fsize,
                color=text_color,
            )

        # Page number (not on cover)
        if idx > 0:
            add_textbox(
                slide, w - Inches(1.2), h - Inches(0.35), Inches(1.0), Inches(0.3),
                str(idx),
                font_size=Pt(11), color=C_ACCENT, align=PP_ALIGN.RIGHT,
            )

    return prs


def main():
    parser = argparse.ArgumentParser(description='Marp MD → PPTX converter')
    parser.add_argument('input', help='Marp Markdown fájl (.md)')
    parser.add_argument('--template', default=None, help='du_template.pptx elérési útja')
    parser.add_argument('--output', default=None, help='Kimeneti .pptx fájl neve')
    args = parser.parse_args()

    md_path = Path(args.input)
    if not md_path.exists():
        sys.exit(f"Nem található: {md_path}")

    out_path = Path(args.output) if args.output else md_path.with_suffix('.pptx')

    with open(md_path, encoding='utf-8') as f:
        md_text = f.read()

    slides_data = parse_marp(md_text)
    print(f"Parsed {len(slides_data)} slides from {md_path.name}")

    prs = build_presentation(slides_data, args.template)
    prs.save(out_path)
    print(f"✅ Mentve: {out_path}")


if __name__ == '__main__':
    main()
